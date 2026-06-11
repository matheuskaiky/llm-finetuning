"""Extraction and exact-deduplication for the docente (SIGAA) corpus.

Walks ``data/raw/docentesDC-sigaa`` (tree ``professor/ano/mes/dia/arquivo``),
triages files into ``text``/``code``/``noise`` buckets, derives metadata from the
path and filename, extracts text per type, deduplicates exact copies (keeping the
most recent version as canonical), and writes one JSONL record per canonical
document.

Heavy/optional parsers (``pypdf``, ``python-docx``, ``python-pptx``) are imported
lazily so the triage, metadata, fingerprint, and dedup logic stay testable without
the ML or document-parsing stack.

Design notes and the closed decisions behind this module live in
``PENSAMENTO_EXTRACAO_DOCENTES.md`` (git-ignored).
"""

from __future__ import annotations

import hashlib
import html
import json
import re
import unicodedata
from collections.abc import Iterable, Iterator, Sequence
from pathlib import Path
from typing import Any

from ..core.interfaces import DatasetLoader
from ..core.registry import DATASET_LOADERS
from .pdf_loader import normalize_text

# Extensions whose body becomes factual text for SFT/RAG.
TEXT_EXTS = {
    ".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".text",
    ".tex", ".html", ".htm", ".csv", ".md", ".rtf",
}
# Source code: excluded from the factual corpus by default, opt-in subcorpus.
CODE_EXTS = {
    ".java", ".py", ".c", ".cpp", ".cc", ".h", ".hpp", ".hs", ".f90", ".f",
    ".js", ".ts", ".sql", ".sh", ".r", ".m", ".go", ".rs", ".rb", ".php",
    ".scala", ".kt", ".cs",
}
# Path components that mark non-authored or already-handled content (always noise).
NOISE_DIRS = {
    "_archives", "__MACOSX", ".venv", "venv", "site-packages",
    "node_modules", ".git", ".ipynb_checkpoints",
}

_SIGAA_PREFIX = re.compile(r"^(\d+)-(.+)$")
_WS = re.compile(r"\s+")
_HTML_TAG = re.compile(r"<[^>]+>")


def classify_file(path: str | Path) -> str:
    """Return the triage bucket of ``path``: ``"text"``, ``"code"`` or ``"noise"``.

    A noise directory anywhere in the path, an AppleDouble stub (``._*``) or a
    dotfile overrides the extension. Otherwise the suffix decides; anything not in
    :data:`TEXT_EXTS` or :data:`CODE_EXTS` is noise.
    """
    path = Path(path)
    parts = path.parts
    if any(part in NOISE_DIRS for part in parts):
        return "noise"
    name = path.name
    if name.startswith("._") or name == ".DS_Store":
        return "noise"
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTS:
        return "text"
    if suffix in CODE_EXTS:
        return "code"
    return "noise"


def parse_sigaa_id(filename: str) -> tuple[str | None, str]:
    """Split a SIGAA numeric id prefix from a filename.

    SIGAA exports prefix files with a numeric id (e.g. ``6228726-Lista.pdf``).
    Returns ``(id, rest)`` when present, else ``(None, filename)``.
    """
    match = _SIGAA_PREFIX.match(filename)
    if match:
        return match.group(1), match.group(2)
    return None, filename


def _as_int(value: str) -> int | None:
    return int(value) if value.isdigit() else None


def parse_path_metadata(root: str | Path, path: str | Path) -> dict[str, Any]:
    """Derive ``professor``/``year``/``month``/``day``/``date`` from the tree.

    The expected layout is ``root/professor/ano/mes/dia/arquivo``. Missing or
    non-numeric date components are left as ``None``; ``date`` is the ISO string of
    whatever parts are present (e.g. ``"2024-03-05"`` or ``"2024"``).
    """
    rel_parts = Path(path).relative_to(root).parts
    professor = rel_parts[0] if rel_parts else None
    year = _as_int(rel_parts[1]) if len(rel_parts) > 1 else None
    month = _as_int(rel_parts[2]) if len(rel_parts) > 2 else None
    day = _as_int(rel_parts[3]) if len(rel_parts) > 3 else None
    date = "-".join(
        f"{v:02d}" if i else f"{v:04d}"
        for i, v in enumerate((year, month, day))
        if v is not None
    )
    return {
        "professor": professor,
        "year": year,
        "month": month,
        "day": day,
        "date": date or None,
    }


def _date_key(record: dict[str, Any]) -> tuple[int, int, int, str]:
    """Sortable key for canonical selection: newest date wins, path breaks ties."""
    date = record.get("date") or ""
    parts = [int(p) if p.isdigit() else 0 for p in date.split("-")]
    parts += [0] * (3 - len(parts))
    return (parts[0], parts[1], parts[2], record.get("source_path", ""))


def text_fingerprint(text: str) -> str:
    """SHA-1 of the text after lowercasing and collapsing all whitespace.

    Captures documents whose binary differs (re-export, changed metadata) but whose
    textual content is identical.
    """
    normalized = _WS.sub(" ", text.lower()).strip()
    return hashlib.sha1(normalized.encode("utf-8")).hexdigest()


def file_md5(path: str | Path, chunk_size: int = 1 << 20) -> str:
    """Streaming MD5 of a file's bytes."""
    digest = hashlib.md5()
    with open(path, "rb") as handle:
        for chunk in iter(lambda: handle.read(chunk_size), b""):
            digest.update(chunk)
    return digest.hexdigest()


def sanitize_text(text: str) -> str:
    """Drop lone surrogates and other non-UTF-8-encodable code points.

    PDF text extraction can emit unpaired surrogate characters that ``json.dumps``
    and UTF-8 cannot encode. Re-encoding with ``errors="ignore"`` removes them while
    keeping all valid (accented) text.
    """
    return text.encode("utf-8", "ignore").decode("utf-8")


def _strip_html(raw: str) -> str:
    """Remove tags and resolve HTML entities (numeric entities cover docx accents)."""
    text = _HTML_TAG.sub(" ", raw)
    return unicodedata.normalize("NFC", html.unescape(text))


def deduplicate(records: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    """Collapse exact duplicates, keeping the most recent version as canonical.

    Records are grouped by ``text_sha1`` when present, else by ``content_md5``. For
    each group the canonical is the newest by ``date`` (ties broken by
    ``source_path``). The canonical carries ``dup_count`` and the sorted
    ``duplicated_dates`` of the group; when the group spans more than one professor
    it also carries ``shared_with_professors`` (shared standard material, not
    attributed to a single docente).
    """
    groups: dict[str, list[dict[str, Any]]] = {}
    for record in records:
        key = record.get("text_sha1") or record.get("content_md5") or ""
        groups.setdefault(key, []).append(record)

    canonicals: list[dict[str, Any]] = []
    for group in groups.values():
        canonical = dict(max(group, key=_date_key))
        canonical["dup_count"] = len(group)
        canonical["duplicated_dates"] = sorted(
            {r["date"] for r in group if r.get("date")}
        )
        professors = sorted({r["professor"] for r in group if r.get("professor")})
        if len(professors) > 1:
            canonical["shared_with_professors"] = professors
        canonicals.append(canonical)

    canonicals.sort(key=lambda r: (r.get("source_path", "")))
    return canonicals


def export_plaintext_corpus(
    records: Iterable[dict[str, Any]],
    out_dir: str | Path,
    min_chars: int = 200,
) -> list[Path]:
    """Materialize canonical documents as one ``.txt`` per doc for pretraining.

    Writes the cleaned ``text`` of each record with at least ``min_chars``
    characters into ``out_dir``, producing a plain-text corpus compatible with
    :class:`~llm_finetuning.data.text_corpus.TextCorpusLoader` (continual
    pretraining). Records below the threshold or without text are skipped to keep
    boilerplate and parser noise out of the corpus.
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []
    for record in records:
        text = record.get("text") or ""
        if len(text) < min_chars:
            continue
        doc_id = record.get("doc_id") or text_fingerprint(text)
        sigaa = record.get("sigaa_id") or "na"
        out_path = out_dir / f"{sigaa}-{doc_id[:12]}.txt"
        out_path.write_text(text, encoding="utf-8")
        written.append(out_path)
    return written


def word_shingles(text: str, k: int = 5) -> set[str]:
    """Return the set of ``k``-word shingles of normalized text.

    Used as the MinHash input for near-duplicate detection. Text shorter than
    ``k`` tokens collapses to a single shingle.
    """
    tokens = _WS.sub(" ", text.lower()).strip().split(" ")
    tokens = [t for t in tokens if t]
    if not tokens:
        return set()
    if len(tokens) < k:
        return {" ".join(tokens)}
    return {" ".join(tokens[i : i + k]) for i in range(len(tokens) - k + 1)}


def _union_find(n: int, edges: Iterable[tuple[int, int]]) -> list[list[int]]:
    """Group ``0..n-1`` into connected components given undirected ``edges``."""
    parent = list(range(n))

    def find(x: int) -> int:
        while parent[x] != x:
            parent[x] = parent[parent[x]]
            x = parent[x]
        return x

    for a, b in edges:
        parent[find(a)] = find(b)
    clusters: dict[int, list[int]] = {}
    for i in range(n):
        clusters.setdefault(find(i), []).append(i)
    return list(clusters.values())


def deduplicate_near(
    records: Sequence[dict[str, Any]],
    threshold: float = 0.85,
    num_perm: int = 128,
    shingle_size: int = 5,
) -> list[dict[str, Any]]:
    """Collapse near-duplicate documents within each professor via MinHash/LSH.

    Runs after the exact pass. Records with text are bucketed by professor; within
    a bucket, documents whose word-shingle Jaccard estimate is at least
    ``threshold`` are clustered (transitively) and collapsed to the most recent
    version, merging ``duplicated_dates`` and summing ``dup_count``. Records
    without text, or alone in their bucket, pass through unchanged. Cross-professor
    overlap is intentionally not merged here (handled as shared material by the
    exact pass).
    """
    from datasketch import MinHash, MinHashLSH

    by_professor: dict[Any, list[int]] = {}
    passthrough: list[dict[str, Any]] = []
    for idx, record in enumerate(records):
        if record.get("text"):
            by_professor.setdefault(record.get("professor"), []).append(idx)
        else:
            passthrough.append(record)

    result: list[dict[str, Any]] = list(passthrough)
    for indices in by_professor.values():
        if len(indices) == 1:
            result.append(records[indices[0]])
            continue
        minhashes: dict[int, Any] = {}
        lsh = MinHashLSH(threshold=threshold, num_perm=num_perm)
        for idx in indices:
            mh = MinHash(num_perm=num_perm, seed=1)
            for shingle in word_shingles(records[idx]["text"], shingle_size):
                mh.update(shingle.encode("utf-8"))
            minhashes[idx] = mh
            lsh.insert(str(idx), mh)
        edges: list[tuple[int, int]] = []
        local = {idx: pos for pos, idx in enumerate(indices)}
        for idx in indices:
            for other in lsh.query(minhashes[idx]):
                j = int(other)
                if j != idx:
                    edges.append((local[idx], local[j]))
        for cluster in _union_find(len(indices), edges):
            group = [records[indices[pos]] for pos in cluster]
            result.append(_collapse_cluster(group))
    result.sort(key=lambda r: r.get("source_path", ""))
    return result


def _collapse_cluster(group: list[dict[str, Any]]) -> dict[str, Any]:
    """Merge a near-duplicate cluster into its most recent member."""
    if len(group) == 1:
        return group[0]
    canonical = dict(max(group, key=_date_key))
    dates: set[str] = set()
    professors: set[str] = set()
    count = 0
    for record in group:
        member_dates = record.get("duplicated_dates")
        if not member_dates:
            member_dates = [record["date"]] if record.get("date") else []
        dates.update(member_dates)
        if record.get("professor"):
            professors.add(record["professor"])
        count += record.get("dup_count", 1)
    canonical["dup_count"] = count
    canonical["duplicated_dates"] = sorted(dates)
    canonical["near_dup_count"] = len(group)
    if len(professors) > 1:
        canonical["shared_with_professors"] = sorted(professors)
    return canonical


@DATASET_LOADERS.register("docente_extractor")
class DocenteExtractor(DatasetLoader):
    """Triage, text-extract and exact-dedup the docente (SIGAA) corpus to JSONL.

    Args:
        input_dir: Root of the docente tree (``professor/ano/mes/dia/arquivo``).
        output_path: JSONL file written with one canonical document per line.
        include_code: When True, source-code files are also extracted as a separate
            subcorpus; by default only the factual text bucket is processed.
        near_dedup: When True, run the near-duplicate pass after the exact one.
        near_threshold: Jaccard threshold for the near-duplicate pass.
        encoding: Encoding used to read plain-text files.
    """

    def __init__(
        self,
        input_dir: str | Path,
        output_path: str | Path,
        include_code: bool = False,
        near_dedup: bool = True,
        near_threshold: float = 0.85,
        encoding: str = "utf-8",
    ) -> None:
        self.input_dir = Path(input_dir)
        self.output_path = Path(output_path)
        self.include_code = include_code
        self.near_dedup = near_dedup
        self.near_threshold = near_threshold
        self.encoding = encoding

    def iter_files(self) -> Iterator[tuple[Path, str]]:
        """Yield ``(path, bucket)`` for every kept file under ``input_dir``.

        Noise is always skipped; code is skipped unless ``include_code`` is set.
        """
        wanted = {"text", "code"} if self.include_code else {"text"}
        for path in sorted(self.input_dir.rglob("*")):
            if not path.is_file():
                continue
            bucket = classify_file(path)
            if bucket in wanted:
                yield path, bucket

    def extract_text(self, path: str | Path) -> tuple[str, str]:
        """Return ``(text, origin)`` for a single file.

        ``origin`` is ``"text_layer"`` for direct/text-layer extraction,
        ``"unavailable"`` when an optional parser is not installed, or
        ``"unsupported"`` for legacy binary formats not handled here.
        """
        path = Path(path)
        suffix = path.suffix.lower()
        if suffix in {".txt", ".text", ".tex", ".csv", ".md"}:
            raw = path.read_text(encoding=self.encoding, errors="replace")
            return normalize_text(raw), "text_layer"
        if suffix in {".html", ".htm"}:
            raw = path.read_text(encoding=self.encoding, errors="replace")
            return normalize_text(_strip_html(raw)), "text_layer"
        if suffix == ".pdf":
            return self._extract_pdf(path), "text_layer"
        if suffix == ".docx":
            return self._extract_docx(path)
        if suffix == ".pptx":
            return self._extract_pptx(path)
        if suffix in CODE_EXTS:
            raw = path.read_text(encoding=self.encoding, errors="replace")
            return normalize_text(raw), "text_layer"
        # Legacy binary office formats need an external converter (LibreOffice).
        return "", "unsupported"

    def _extract_pdf(self, path: Path) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return normalize_text("\n\n".join(pages))

    def _extract_docx(self, path: Path) -> tuple[str, str]:
        try:
            import docx  # type: ignore
        except ImportError:
            return "", "unavailable"
        document = docx.Document(str(path))
        text = "\n".join(p.text for p in document.paragraphs)
        return normalize_text(_strip_html(text)), "text_layer"

    def _extract_pptx(self, path: Path) -> tuple[str, str]:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return "", "unavailable"
        chunks: list[str] = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
        return normalize_text("\n".join(chunks)), "text_layer"

    def _build_record(self, path: Path, bucket: str) -> dict[str, Any]:
        try:
            text, origin = self.extract_text(path)
        except Exception:
            # A single unreadable/corrupt file must not abort the whole run; it is
            # kept as a metadata-only record so the dedup and audit still see it.
            text, origin = "", "error"
        text = sanitize_text(text)
        meta = parse_path_metadata(self.input_dir, path)
        sigaa_id, clean_name = parse_sigaa_id(path.name)
        text_sha1 = text_fingerprint(text) if text else ""
        content_md5 = file_md5(path)
        return {
            "doc_id": text_sha1 or content_md5,
            "professor": meta["professor"],
            "date": meta["date"],
            "year": meta["year"],
            "month": meta["month"],
            "day": meta["day"],
            "sigaa_id": sigaa_id,
            "source_path": str(path.relative_to(self.input_dir)),
            "filename_clean": clean_name,
            "filetype": path.suffix.lower().lstrip("."),
            "bucket": bucket,
            "text_origin": origin,
            "content_md5": content_md5,
            "text_sha1": text_sha1,
            "text": text,
        }

    def load(self) -> list[dict[str, Any]]:
        """Extract, deduplicate and write the JSONL; return the canonical records."""
        raw_records = [
            self._build_record(path, bucket) for path, bucket in self.iter_files()
        ]
        records = deduplicate(raw_records)
        if self.near_dedup:
            records = deduplicate_near(records, threshold=self.near_threshold)
        self._write_jsonl(records)
        return records

    def _write_jsonl(self, records: Sequence[dict[str, Any]]) -> None:
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.output_path, "w", encoding="utf-8") as handle:
            for record in records:
                # Sanitize the serialized line as a whole: path-derived fields
                # (source_path, filename_clean, professor) can carry lone
                # surrogates from undecodable filesystem bytes, not just ``text``.
                line = sanitize_text(json.dumps(record, ensure_ascii=False))
                handle.write(line + "\n")
